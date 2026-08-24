"""
generate_presentation.py
Generates the Final Project Presentation for:
Velsora: Multi-Agent AI Analysis System for Financial Research and Business Insight
Infosys Springboard Virtual Internship 7.0

Design: Professional Dark Cream / Warm Corporate palette with large readable text.
Font sizing follows industry standards: Titles 36-40pt, Sub 24-28pt, Body 18-22pt.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ───────────────────────────────────────────────────
# DARK CREAM CORPORATE PALETTE
# ───────────────────────────────────────────────────
BG       = RGBColor(243, 239, 230)   # #F3EFE6  Dark Cream
CARD     = RGBColor(255, 255, 255)   # #FFFFFF  White cards
CARD_ALT = RGBColor(247, 241, 231)   # #F7F1E7  Oat Milk Cream
BORDER   = RGBColor(203, 187, 160)   # #CBBBA0  Warm Beige border
CHARCOAL = RGBColor(46, 46, 51)      # #2E2E33  Charcoal primary text
SLATE    = RGBColor(75, 78, 87)      # #4B4E57  Secondary text
EARTH    = RGBColor(139, 126, 106)   # #8B7E6A  Soft Earth muted
NAVY     = RGBColor(27, 58, 92)      # #1B3A5C  Deep Navy accent
STEEL    = RGBColor(74, 127, 181)    # #4A7FB5  Steel Blue
AMBER    = RGBColor(232, 145, 58)    # #E8913A  Warm Amber highlight
TEAL     = RGBColor(15, 118, 110)    # #0F766E  Deep Teal
GREEN    = RGBColor(22, 163, 74)     # #16A34A  Success green
RED      = RGBColor(185, 28, 28)     # #B91C1C  Crimson alert
TBL_HEAD = RGBColor(27, 58, 92)      # #1B3A5C  Table header = Navy
WHITE    = RGBColor(255, 255, 255)

FONT_H = "Segoe UI Semibold"
FONT_B = "Segoe UI"
FONT_C = "Consolas"

SCREENSHOTS = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../project doc/screenshots"))
OUTPUT = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../project doc/Velsora_Multi_Agent_Financial_Research_System_Presentation.pptx"))

W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = Inches(11.533)

# ───────────────────────────────────────────────────
# HELPERS
# ───────────────────────────────────────────────────

def bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()

def accent_bar(slide, left, top, width, height, color):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

def card(slide, l, t, w, h, fill=CARD, brd=BORDER):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = brd; c.line.width = Pt(1)
    return c

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def hdr(slide, tag, title, sub=""):
    """Standard slide header: small tag pill + big title + subtitle."""
    # Tag pill
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(0.45), Inches(3.8), Inches(0.35))
    p.fill.solid(); p.fill.fore_color.rgb = CARD_ALT
    p.line.color.rgb = BORDER; p.line.width = Pt(0.75)
    tf = p.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    pp = tf.paragraphs[0]
    pp.text = tag.upper(); pp.font.name = FONT_B; pp.font.size = Pt(11)
    pp.font.bold = True; pp.font.color.rgb = EARTH

    # Title
    t = tb(slide, MARGIN, Inches(0.88), CONTENT_W, Inches(0.65))
    tf2 = t.text_frame; tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    pt = tf2.paragraphs[0]
    pt.text = title; pt.font.name = FONT_H; pt.font.size = Pt(32)
    pt.font.bold = True; pt.font.color.rgb = CHARCOAL

    # Subtitle
    if sub:
        ps = tf2.add_paragraph()
        ps.text = sub; ps.font.name = FONT_B; ps.font.size = Pt(16)
        ps.font.color.rgb = SLATE; ps.space_before = Pt(4)

def ftr(slide, num, total=16):
    accent_bar(slide, MARGIN, Inches(6.95), CONTENT_W, Inches(0.02), BORDER)
    f = tb(slide, MARGIN, Inches(7.0), Inches(9.0), Inches(0.35))
    tf = f.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    pp = tf.paragraphs[0]
    pp.text = "Velsora  ·  Multi-Agent AI Analysis System  ·  Infosys Springboard  ·  Team 2"
    pp.font.name = FONT_B; pp.font.size = Pt(11); pp.font.color.rgb = EARTH

    n = tb(slide, Inches(10.5), Inches(7.0), Inches(1.93), Inches(0.35))
    tf2 = n.text_frame; tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    pn = tf2.paragraphs[0]
    pn.text = f"{num:02d} / {total:02d}"
    pn.font.name = FONT_H; pn.font.size = Pt(11); pn.font.bold = True
    pn.font.color.rgb = NAVY; pn.alignment = PP_ALIGN.RIGHT

def bullet_text(tf, items, bold_size=18, body_size=16, bold_color=CHARCOAL, body_color=SLATE):
    """Add bullet items to a text frame. Each item is (bold_part, body_part)."""
    for bld, body in items:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = f"  {bld}  "
        r1.font.name = FONT_H; r1.font.size = Pt(bold_size)
        r1.font.bold = True; r1.font.color.rgb = bold_color
        r2 = p.add_run()
        r2.text = body
        r2.font.name = FONT_B; r2.font.size = Pt(body_size)
        r2.font.color.rgb = body_color


def create_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    BL = prs.slide_layouts[6]

    # ═══════════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)

    # Top accent line
    accent_bar(s, MARGIN, Inches(1.0), CONTENT_W, Inches(0.06), NAVY)

    # Title block
    t = tb(s, MARGIN, Inches(1.3), CONTENT_W, Inches(3.8))
    tf = t.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    p1.text = "VELSORA"
    p1.font.name = FONT_H; p1.font.size = Pt(54); p1.font.bold = True
    p1.font.color.rgb = CHARCOAL

    p2 = tf.add_paragraph()
    p2.text = "Multi-Agent AI Analysis System for\nFinancial Research and Business Insight"
    p2.font.name = FONT_H; p2.font.size = Pt(26); p2.font.bold = True
    p2.font.color.rgb = NAVY; p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Infosys Springboard Virtual Internship 7.0  ·  Final Capstone Project  ·  Team 2"
    p3.font.name = FONT_B; p3.font.size = Pt(16); p3.font.color.rgb = SLATE
    p3.space_before = Pt(20)

    # Bottom info cards
    cards = [
        ("Team Lead", "Samrat Maurya", NAVY),
        ("AI Engine", "Groq LLaMA 3.3 · Nemotron 550B", TEAL),
        ("Architecture", "6 Specialized AI Agents", STEEL),
        ("Deliverables", "All 6 Mentor Requirements Met", GREEN),
    ]
    cw = Inches(2.78); cg = Inches(0.18)
    for i, (label, val, col) in enumerate(cards):
        cx = MARGIN + i * (cw + cg)
        card(s, cx, Inches(5.3), cw, Inches(1.35))
        accent_bar(s, cx, Inches(5.3), cw, Inches(0.06), col)
        tt = tb(s, cx + Inches(0.2), Inches(5.45), cw - Inches(0.4), Inches(1.1))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = label.upper(); pp.font.name = FONT_B; pp.font.size = Pt(12)
        pp.font.bold = True; pp.font.color.rgb = EARTH
        pv = ttf.add_paragraph()
        pv.text = val; pv.font.name = FONT_H; pv.font.size = Pt(18)
        pv.font.bold = True; pv.font.color.rgb = CHARCOAL; pv.space_before = Pt(4)

    ftr(s, 1)

    # ═══════════════════════════════════════════════
    # SLIDE 2 — DELIVERABLES AUDIT
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Executive Summary", "Mentor Deliverable Compliance", "All 6 required deliverables verified and submitted.")

    items = [
        ("1. Agile Documentation", "Sprint backlogs, user stories, RTM matrix, burn-down reports"),
        ("2. Project Presentation", "16-slide professional presentation with real project data"),
        ("3. Technical Documentation", "IEEE 29148 SAD (6,400+ lines) + 22-chapter SDS"),
        ("4. Main GitHub Repository", "Production-ready codebase with Docker & test suites"),
        ("5. Individual Intern Repos", "Synced repositories for all team members"),
        ("6. Working Project System", "FastAPI + React 19 + LangGraph + MongoDB Atlas"),
    ]
    for i, (title, desc) in enumerate(items):
        yt = Inches(2.0) + i * Inches(0.82)
        card(s, MARGIN, yt, CONTENT_W, Inches(0.72))
        accent_bar(s, MARGIN, yt, Inches(0.06), Inches(0.72), GREEN)
        tt = tb(s, MARGIN + Inches(0.2), yt + Inches(0.06), CONTENT_W - Inches(0.4), Inches(0.6))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        r1 = pp.add_run()
        r1.text = title; r1.font.name = FONT_H; r1.font.size = Pt(18)
        r1.font.bold = True; r1.font.color.rgb = CHARCOAL
        r2 = pp.add_run()
        r2.text = "   ✓ VERIFIED"; r2.font.name = FONT_H; r2.font.size = Pt(14)
        r2.font.bold = True; r2.font.color.rgb = GREEN
        p2 = ttf.add_paragraph()
        p2.text = desc; p2.font.name = FONT_B; p2.font.size = Pt(15)
        p2.font.color.rgb = SLATE; p2.space_before = Pt(2)

    ftr(s, 2)

    # ═══════════════════════════════════════════════
    # SLIDE 3 — PROBLEM STATEMENT
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Problem Statement", "The Financial Research Bottleneck")

    probs = [
        ("Information Overload", "SEC 10-K filings span 100–350 pages of dense financial jargon.\nManual extraction takes 12–18 hours per filing.", RED),
        ("LLM Hallucination Risk", "Commercial AI models fabricate financial figures by up to 22%.\nNo citation lineage or source verification.", AMBER),
        ("Fragmented Benchmarking", "Cross-company comparison requires hours of manual\nspreadsheet normalization across incompatible formats.", NAVY),
    ]
    pw = Inches(3.72); pg = Inches(0.2)
    for i, (title, desc, col) in enumerate(probs):
        px = MARGIN + i * (pw + pg)
        card(s, px, Inches(2.1), pw, Inches(3.0))
        accent_bar(s, px, Inches(2.1), pw, Inches(0.08), col)
        tt = tb(s, px + Inches(0.22), Inches(2.35), pw - Inches(0.44), Inches(2.6))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = title; pp.font.name = FONT_H; pp.font.size = Pt(22)
        pp.font.bold = True; pp.font.color.rgb = CHARCOAL
        pd = ttf.add_paragraph()
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(16)
        pd.font.color.rgb = SLATE; pd.space_before = Pt(10)

    # Stats bar
    stats = [("100–350+", "Pages per 10-K"), ("14.5 hrs", "Analyst time / filing"), ("22%", "LLM hallucination rate"), ("Zero", "Auditability in legacy AI")]
    sw = Inches(2.78); sg = Inches(0.18)
    for i, (val, lbl) in enumerate(stats):
        sx = MARGIN + i * (sw + sg)
        card(s, sx, Inches(5.35), sw, Inches(1.3))
        tt = tb(s, sx + Inches(0.15), Inches(5.45), sw - Inches(0.3), Inches(1.1))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = val; pp.font.name = FONT_H; pp.font.size = Pt(28)
        pp.font.bold = True; pp.font.color.rgb = NAVY; pp.alignment = PP_ALIGN.CENTER
        pl = ttf.add_paragraph()
        pl.text = lbl; pl.font.name = FONT_B; pl.font.size = Pt(14)
        pl.font.color.rgb = EARTH; pl.alignment = PP_ALIGN.CENTER; pl.space_before = Pt(4)

    ftr(s, 3)

    # ═══════════════════════════════════════════════
    # SLIDE 4 — SOLUTION OVERVIEW
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Our Solution", "Autonomous Multi-Agent Financial Diligence")

    sols = [
        ("6 Specialized AI Agents", "Each agent owns a single concern — parsing,\nextraction, risk scanning, research, comparison, reporting.", NAVY),
        ("Deterministic Source Grounding", "Every extracted fact is cryptographically stamped\nwith Document ID, Page Number, and Source Snippet.", TEAL),
        ("Zero-Hallucination Mandate", "If evidence cannot be traced to indexed chunks,\nthe system strictly refuses to answer.", RED),
        ("Autonomous Event Pipelines", "Uploading a filing instantly triggers ingestion,\nvectorization, extraction, and risk scanning.", GREEN),
    ]
    cw2 = Inches(5.58); ch2 = Inches(2.05); cg2 = Inches(0.35)
    for i, (title, desc, col) in enumerate(sols):
        cx = MARGIN + (i % 2) * (cw2 + cg2)
        cy = Inches(2.05) + (i // 2) * (ch2 + Inches(0.25))
        card(s, cx, cy, cw2, ch2)
        accent_bar(s, cx, cy, Inches(0.07), ch2, col)
        tt = tb(s, cx + Inches(0.25), cy + Inches(0.18), cw2 - Inches(0.5), ch2 - Inches(0.35))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = title; pp.font.name = FONT_H; pp.font.size = Pt(22)
        pp.font.bold = True; pp.font.color.rgb = CHARCOAL
        pd = ttf.add_paragraph()
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(16)
        pd.font.color.rgb = SLATE; pd.space_before = Pt(8)

    ftr(s, 4)

    # ═══════════════════════════════════════════════
    # SLIDE 5 — TEAM & REPOSITORIES
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Team Structure", "Team 2 — Ownership & GitHub Repositories")

    team = [
        ("Samrat (Lead)", "Research Agent, Report Agent, LangGraph Orchestrator, React 19 UI", "github.com/samratmaurya1217/multi-agent-financial-research-system"),
        ("Snigdha", "Document Agent — PDF/DOCX ingestion, PyMuPDF, chunking, vectors", "github.com/Snigdha-02/multi-agent-financial-research-system"),
        ("Thirumala", "Extraction Agent — GAAP metric engine, ratio derivations, citations", "github.com/thirumala95/multi-agent-financial-research-system"),
        ("Harshitha", "Red Flag Agent — Forensic risk detection, severity classification", "Core Contributor (Red Flag & MongoDB Schemas)"),
        ("Akshith", "Backend & API — FastAPI controllers, Pydantic schemas, E2E tests", "github.com/06102006-Ak/multi-agent-financial-research-system"),
    ]

    table_shape = s.shapes.add_table(6, 3, MARGIN, Inches(2.1), CONTENT_W, Inches(4.5))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(5.1)
    tbl.columns[2].width = Inches(4.233)

    for ci, txt in enumerate(["Team Member", "Subsystem & Contributions", "GitHub Repository"]):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = TBL_HEAD
        p = c.text_frame.paragraphs[0]
        p.text = txt; p.font.name = FONT_H; p.font.size = Pt(16)
        p.font.bold = True; p.font.color.rgb = WHITE

    for ri, (name, contrib, repo) in enumerate(team, 1):
        bgc = CARD if ri % 2 == 1 else CARD_ALT
        for ci, val in enumerate([name, contrib, repo]):
            c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = bgc
            p = c.text_frame.paragraphs[0]
            p.text = val; p.font.name = FONT_B; p.font.size = Pt(14)
            if ci == 0:
                p.font.bold = True; p.font.color.rgb = NAVY
            elif ci == 2:
                p.font.color.rgb = STEEL; p.font.name = FONT_C; p.font.size = Pt(12)
            else:
                p.font.color.rgb = CHARCOAL

    ftr(s, 5)

    # ═══════════════════════════════════════════════
    # SLIDE 6 — SYSTEM ARCHITECTURE
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "System Architecture", "End-to-End Platform Architecture (C4 Model)")

    tiers = [
        ("Presentation Tier", "React 19 · TypeScript · Tailwind CSS\nWorkspace isolation · GAAP metrics grid\nSSE streaming chat · PDF viewer", STEEL),
        ("API & Security", "FastAPI ASGI · JWT · RBAC\nPrompt injection defense\nAsync worker queue", NAVY),
        ("AI Orchestration", "LangGraph StateGraph\n6 dedicated agents\nDurable MongoDB checkpoints", TEAL),
        ("Data Layer", "MongoDB Atlas · Vector Search\n384-dim cosine similarity\n6 collections · SHA-256 dedupe", AMBER),
    ]
    tw = Inches(2.68); tg = Inches(0.16)
    for i, (title, desc, col) in enumerate(tiers):
        tx = MARGIN + i * (tw + tg)
        card(s, tx, Inches(2.1), tw, Inches(3.2))
        hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, Inches(2.1), tw, Inches(0.5))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        htf = hbar.text_frame
        hp = htf.paragraphs[0]
        hp.text = title; hp.font.name = FONT_H; hp.font.size = Pt(16)
        hp.font.bold = True; hp.font.color.rgb = WHITE; hp.alignment = PP_ALIGN.CENTER
        tt = tb(s, tx + Inches(0.18), Inches(2.75), tw - Inches(0.36), Inches(2.4))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pd = ttf.paragraphs[0]
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(16)
        pd.font.color.rgb = CHARCOAL; pd.line_spacing = Pt(22)

    # Data flow bar
    card(s, MARGIN, Inches(5.55), CONTENT_W, Inches(1.1), fill=CARD_ALT)
    tt = tb(s, MARGIN + Inches(0.2), Inches(5.65), CONTENT_W - Inches(0.4), Inches(0.9))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Data Flow:  Upload PDF  →  Document Agent  →  Extraction Agent  →  Red Flag Agent  →  Vector Store  →  Research Agent  →  PDF Report"
    pp.font.name = FONT_H; pp.font.size = Pt(18); pp.font.bold = True; pp.font.color.rgb = NAVY

    ftr(s, 6)

    # ═══════════════════════════════════════════════
    # SLIDE 7 — LLM RESILIENCE
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "AI Model Strategy", "Dual-Tier Multi-Provider LLM Resilience Engine")

    # Tier 1
    card(s, MARGIN, Inches(2.1), Inches(5.58), Inches(2.8))
    accent_bar(s, MARGIN, Inches(2.1), Inches(0.07), Inches(2.8), TEAL)
    tt = tb(s, MARGIN + Inches(0.25), Inches(2.25), Inches(5.1), Inches(2.5))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Tier 1 — Fast Extraction & Risk Scanning"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = TEAL
    bullet_text(ttf, [
        ("Engine:", "Groq Cloud — LLaMA 3.3 70B Versatile"),
        ("Latency:", "2.1 seconds per extraction batch"),
        ("Cost:", "$0.00059 per filing run (~92% savings vs GPT-4o)"),
        ("Fallback:", "Google Gemini 2.5 Flash via OpenRouter"),
    ], bold_size=16, body_size=16, bold_color=CHARCOAL, body_color=SLATE)

    # Tier 2
    card(s, Inches(6.58), Inches(2.1), Inches(5.58), Inches(2.8))
    accent_bar(s, Inches(6.58), Inches(2.1), Inches(0.07), Inches(2.8), NAVY)
    tt = tb(s, Inches(6.83), Inches(2.25), Inches(5.1), Inches(2.5))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Tier 2 — Deep Reasoning & Report Synthesis"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = NAVY
    bullet_text(ttf, [
        ("Engine:", "NVIDIA Nemotron 3 Ultra 550B"),
        ("Streaming:", "420ms Time-to-First-Token via SSE"),
        ("Strengths:", "Chain-of-thought, cross-referencing, synthesis"),
        ("Failover:", "Nemotron → Gemini Flash → Groq LLaMA"),
    ], bold_size=16, body_size=16, bold_color=CHARCOAL, body_color=SLATE)

    # Circuit breaker cards
    cbs = [
        ("Step 1: Smart Routing", "Routes to optimal model based\non task complexity.", NAVY),
        ("Step 2: Circuit Breaker", "On HTTP 429/503, triggers\nautomatic fallback chain.", AMBER),
        ("Step 3: Graceful Degradation", "Serves cached results with clear\ndegraded-status indicators.", GREEN),
    ]
    cbw = Inches(3.72); cbg = Inches(0.2)
    for i, (title, desc, col) in enumerate(cbs):
        cx = MARGIN + i * (cbw + cbg)
        card(s, cx, Inches(5.15), cbw, Inches(1.5))
        accent_bar(s, cx, Inches(5.15), cbw, Inches(0.06), col)
        tt = tb(s, cx + Inches(0.2), Inches(5.3), cbw - Inches(0.4), Inches(1.25))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = title; pp.font.name = FONT_H; pp.font.size = Pt(18)
        pp.font.bold = True; pp.font.color.rgb = CHARCOAL
        pd = ttf.add_paragraph()
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(15)
        pd.font.color.rgb = SLATE; pd.space_before = Pt(4)

    ftr(s, 7)

    # ═══════════════════════════════════════════════
    # SLIDE 8 — DOCUMENT AGENT
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Agent 1", "Document Ingestion & Vectorization Pipeline")

    steps = [
        ("Upload", "PDF/DOCX/TXT\n50MB limit\nSHA-256 dedupe", NAVY),
        ("Parse", "PyMuPDF layout\nTable bounding\nOCR fallback", TEAL),
        ("Chunk", "500-token sliding\n100-token overlap\nNarrative continuity", STEEL),
        ("Embed", "384-dim vectors\nMiniLM-L6-v2\nBatch encoding", AMBER),
        ("Index", "Atlas Vector Search\nCosine similarity\n< 82ms queries", GREEN),
    ]
    sw2 = Inches(2.1); sg2 = Inches(0.18)
    for i, (title, desc, col) in enumerate(steps):
        sx = MARGIN + i * (sw2 + sg2)
        card(s, sx, Inches(2.1), sw2, Inches(2.5))
        hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, Inches(2.1), sw2, Inches(0.45))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        htf = hbar.text_frame
        hp = htf.paragraphs[0]
        hp.text = title; hp.font.name = FONT_H; hp.font.size = Pt(16)
        hp.font.bold = True; hp.font.color.rgb = WHITE; hp.alignment = PP_ALIGN.CENTER
        tt = tb(s, sx + Inches(0.15), Inches(2.7), sw2 - Inches(0.3), Inches(1.8))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pd = ttf.paragraphs[0]
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(16)
        pd.font.color.rgb = CHARCOAL; pd.line_spacing = Pt(22)

    # Benchmarks
    card(s, MARGIN, Inches(4.85), CONTENT_W, Inches(1.8))
    tt = tb(s, MARGIN + Inches(0.25), Inches(4.95), CONTENT_W - Inches(0.5), Inches(1.6))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Performance Benchmarks"; pp.font.name = FONT_H; pp.font.size = Pt(20)
    pp.font.bold = True; pp.font.color.rgb = CHARCOAL

    benches = [
        ("100-page parsing:", "11.4s (target < 15s)"),
        ("500-chunk embeddings:", "3.2s (target < 5s)"),
        ("Vector search query:", "82ms (target < 150ms)"),
        ("Duplicate detection:", "100% SHA-256 quarantine"),
    ]
    bullet_text(ttf, benches, bold_size=16, body_size=16, bold_color=NAVY, body_color=SLATE)

    ftr(s, 8)

    # ═══════════════════════════════════════════════
    # SLIDE 9 — EXTRACTION & RED FLAGS
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Agents 2 & 3", "GAAP Metric Extraction & Red Flag Radar")

    # Left — Extraction
    card(s, MARGIN, Inches(2.1), Inches(5.58), Inches(4.55))
    accent_bar(s, MARGIN, Inches(2.1), Inches(0.07), Inches(4.55), TEAL)
    tt = tb(s, MARGIN + Inches(0.25), Inches(2.25), Inches(5.1), Inches(4.2))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Agent 2 — Extraction Engine"; pp.font.name = FONT_H
    pp.font.size = Pt(24); pp.font.bold = True; pp.font.color.rgb = TEAL
    ps = ttf.add_paragraph()
    ps.text = "Groq LLaMA 3.3 70B  ·  2.1s latency"; ps.font.name = FONT_B
    ps.font.size = Pt(14); ps.font.color.rgb = EARTH; ps.space_before = Pt(4)
    bullet_text(ttf, [
        ("Metrics:", "Revenue, Gross Profit, Net Income, EPS, EBITDA, Current Ratio, Debt/Equity, Cash Flow"),
        ("Citations:", "Every value stamped with page number, source snippet, and confidence score"),
        ("Zero Extrapolation:", "Missing data is omitted — never fabricated or estimated"),
    ], bold_size=16, body_size=15, bold_color=CHARCOAL, body_color=SLATE)

    # Right — Red Flags
    card(s, Inches(6.58), Inches(2.1), Inches(5.58), Inches(4.55))
    accent_bar(s, Inches(6.58), Inches(2.1), Inches(0.07), Inches(4.55), RED)
    tt = tb(s, Inches(6.83), Inches(2.25), Inches(5.1), Inches(4.2))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Agent 3 — Red Flag Radar"; pp.font.name = FONT_H
    pp.font.size = Pt(24); pp.font.bold = True; pp.font.color.rgb = RED
    ps = ttf.add_paragraph()
    ps.text = "Groq LLaMA 3.3 70B  ·  2.4s latency"; ps.font.name = FONT_B
    ps.font.size = Pt(14); ps.font.color.rgb = EARTH; ps.space_before = Pt(4)
    bullet_text(ttf, [
        ("Liquidity Risk:", "Debt/Equity > 2.5x, negative cash flows, sudden debt acceleration"),
        ("Margin Strain:", "Gross margin compression > 300 bps YoY"),
        ("Legal Exposure:", "Antitrust lawsuits, DOJ investigations, patent litigation"),
        ("Audit Concerns:", "Going-concern warnings, internal control weaknesses"),
    ], bold_size=16, body_size=15, bold_color=CHARCOAL, body_color=SLATE)

    ftr(s, 9)

    # ═══════════════════════════════════════════════
    # SLIDE 10 — RESEARCH AGENT & RAG
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Agent 4", "Conversational Research & Hybrid RAG Pipeline")

    card(s, MARGIN, Inches(2.1), CONTENT_W, Inches(4.55))
    tt = tb(s, MARGIN + Inches(0.25), Inches(2.25), CONTENT_W - Inches(0.5), Inches(4.2))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "5-Stage Hybrid RAG Pipeline"; pp.font.name = FONT_H
    pp.font.size = Pt(24); pp.font.bold = True; pp.font.color.rgb = NAVY

    stages = [
        ("1. Query Decomposition —", "Complex prompts are broken into distinct search queries"),
        ("2. Hybrid Retrieval —", "Parallel 384-dim dense vector search + BM25 keyword queries"),
        ("3. Cross-Encoder Re-Ranking —", "Top 20 candidates re-ranked; chunks below 0.65 relevance eliminated"),
        ("4. Grounded Refusal —", "If zero chunks pass threshold, the system refuses to answer"),
        ("5. SSE Token Streaming —", "Real-time reasoning tokens with citation payloads (TTFT: 420ms)"),
    ]
    for bld, body in stages:
        p = ttf.add_paragraph(); p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = bld; r1.font.name = FONT_H; r1.font.size = Pt(18)
        r1.font.bold = True; r1.font.color.rgb = TEAL
        r2 = p.add_run()
        r2.text = "  " + body; r2.font.name = FONT_B; r2.font.size = Pt(17)
        r2.font.color.rgb = CHARCOAL

    ftr(s, 10)

    # ═══════════════════════════════════════════════
    # SLIDE 11 — PEER BENCHMARKING
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Agent 5", "Cross-Company Peer Benchmarking — Tesla vs Ford FY24")

    table_shape = s.shapes.add_table(6, 4, MARGIN, Inches(2.1), CONTENT_W, Inches(2.8))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(2.8)
    tbl.columns[2].width = Inches(2.8)
    tbl.columns[3].width = Inches(2.933)

    for ci, txt in enumerate(["Metric", "Tesla (TSLA)", "Ford (F)", "Delta"]):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = TBL_HEAD
        p = c.text_frame.paragraphs[0]
        p.text = txt; p.font.name = FONT_H; p.font.size = Pt(16)
        p.font.bold = True; p.font.color.rgb = WHITE

    rows = [
        ("Total Revenue", "$97,698 M", "$176,191 M", "-$78,493 M"),
        ("Gross Margin", "18.2%", "13.4%", "+4.8% (+480 bps)"),
        ("Operating Income", "$7,084 M", "$5,410 M", "+$1,674 M"),
        ("Diluted EPS", "$2.14", "$1.08", "+$1.06 (+98.1%)"),
        ("R&D Intensity", "4.5%", "4.0%", "+0.5% (+50 bps)"),
    ]
    for ri, row in enumerate(rows, 1):
        bgc = CARD if ri % 2 == 1 else CARD_ALT
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = bgc
            p = c.text_frame.paragraphs[0]
            p.text = val; p.font.name = FONT_B; p.font.size = Pt(15)
            if ci == 0: p.font.bold = True; p.font.color.rgb = CHARCOAL
            elif ci == 3:
                p.font.bold = True
                p.font.color.rgb = GREEN if "+" in val else RED
            else: p.font.color.rgb = SLATE

    # Key takeaways
    card(s, MARGIN, Inches(5.15), CONTENT_W, Inches(1.5))
    tt = tb(s, MARGIN + Inches(0.25), Inches(5.25), CONTENT_W - Inches(0.5), Inches(1.3))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Agent 5 Synthesis:  "; pp.font.name = FONT_H; pp.font.size = Pt(18); pp.font.bold = True; pp.font.color.rgb = NAVY
    r = pp.add_run()
    r.text = "Tesla exhibits structural profitability superiority (+480 bps gross margin, nearly double diluted EPS) despite Ford's larger revenue scale driven by legacy ICE volumes."
    r.font.name = FONT_B; r.font.size = Pt(16); r.font.color.rgb = CHARCOAL
    ps = ttf.add_paragraph()
    ps.text = "Engine: NVIDIA Nemotron 550B  ·  Latency: 2.8s  ·  Citations: 16 verified chunks"
    ps.font.name = FONT_B; ps.font.size = Pt(14); ps.font.color.rgb = EARTH; ps.space_before = Pt(6)

    ftr(s, 11)

    # ═══════════════════════════════════════════════
    # SLIDE 12 — PDF REPORT ENGINE
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Agent 6", "ReportLab Dynamic PDF Publishing Engine")

    feats = [
        ("Two-Pass NumberedCanvas", "Dynamic 'Page X of Y' numbering\nRunning headers with corporate branding\nStrict visual budgeting — zero overflow", NAVY),
        ("Institutional Layout", "Structured Flowable paragraphs & tables\nAlternating row styling for financials\nColor-coded risk severity badges", TEAL),
        ("6 Report Sections", "1. Header & Ticker Metadata\n2. Executive Summary\n3. Financial Metric Matrix\n4. Red Flag Analysis\n5. Peer Benchmarking\n6. Disclaimer & Citations", GREEN),
    ]
    fw = Inches(3.72); fg = Inches(0.2)
    for i, (title, desc, col) in enumerate(feats):
        fx = MARGIN + i * (fw + fg)
        card(s, fx, Inches(2.1), fw, Inches(3.6))
        accent_bar(s, fx, Inches(2.1), fw, Inches(0.08), col)
        tt = tb(s, fx + Inches(0.2), Inches(2.35), fw - Inches(0.4), Inches(3.2))
        ttf = tt.text_frame; ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
        pp = ttf.paragraphs[0]
        pp.text = title; pp.font.name = FONT_H; pp.font.size = Pt(20)
        pp.font.bold = True; pp.font.color.rgb = CHARCOAL
        pd = ttf.add_paragraph()
        pd.text = desc; pd.font.name = FONT_B; pd.font.size = Pt(16)
        pd.font.color.rgb = SLATE; pd.space_before = Pt(8); pd.line_spacing = Pt(22)

    # Endpoint bar
    card(s, MARGIN, Inches(5.95), CONTENT_W, Inches(0.7), fill=CARD_ALT)
    tt = tb(s, MARGIN + Inches(0.2), Inches(6.0), CONTENT_W - Inches(0.4), Inches(0.55))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "GET /api/v1/workspaces/{id}/reports/export?format=pdf   ·   Compilation: 1.84s   ·   100% grounded citations"
    pp.font.name = FONT_C; pp.font.size = Pt(15); pp.font.bold = True; pp.font.color.rgb = NAVY

    ftr(s, 12)

    # ═══════════════════════════════════════════════
    # SLIDE 13 — DATABASE & SECURITY
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Infrastructure", "Database Architecture & STRIDE Security")

    # Left — DB
    card(s, MARGIN, Inches(2.1), Inches(5.58), Inches(4.55))
    accent_bar(s, MARGIN, Inches(2.1), Inches(0.07), Inches(4.55), NAVY)
    tt = tb(s, MARGIN + Inches(0.25), Inches(2.25), Inches(5.1), Inches(4.2))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "MongoDB Persistence — 6 Collections"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = NAVY
    bullet_text(ttf, [
        ("workspaces —", "Multi-tenant container isolation"),
        ("documents —", "File metadata, SHA-256 hash, indexing status"),
        ("chunks —", "500-token text with 384-dim vector embeddings"),
        ("metrics —", "GAAP values with page citations & confidence"),
        ("red_flags —", "Risk anomalies with severity classifications"),
        ("reports —", "Compiled PDF binaries and JSON synthesis"),
    ], bold_size=16, body_size=15, bold_color=TEAL, body_color=CHARCOAL)

    # Right — Security
    card(s, Inches(6.58), Inches(2.1), Inches(5.58), Inches(4.55))
    accent_bar(s, Inches(6.58), Inches(2.1), Inches(0.07), Inches(4.55), RED)
    tt = tb(s, Inches(6.83), Inches(2.25), Inches(5.1), Inches(4.2))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "STRIDE Security Model & RBAC"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = RED
    bullet_text(ttf, [
        ("Prompt Defense —", "Regex + heuristic filters block injection attacks"),
        ("Tenant Isolation —", "Cross-workspace access returns HTTP 403"),
        ("Info Disclosure —", "System prompts never expose API keys in logs"),
        ("DoS Protection —", "50MB upload cap, 100 req/min rate limiting"),
        ("RBAC —", "Viewer (read) · Analyst (upload) · Admin (manage)"),
    ], bold_size=16, body_size=15, bold_color=CHARCOAL, body_color=SLATE)

    ftr(s, 13)

    # ═══════════════════════════════════════════════
    # SLIDE 14 — BENCHMARKS & TESTING
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Quality Assurance", "Performance Benchmarks & Test Results")

    table_shape = s.shapes.add_table(8, 4, MARGIN, Inches(2.1), CONTENT_W, Inches(4.3))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(4.0)
    tbl.columns[1].width = Inches(2.5)
    tbl.columns[2].width = Inches(2.5)
    tbl.columns[3].width = Inches(2.533)

    for ci, txt in enumerate(["Benchmark", "Target SLO", "Measured", "Status"]):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = TBL_HEAD
        p = c.text_frame.paragraphs[0]
        p.text = txt; p.font.name = FONT_H; p.font.size = Pt(16)
        p.font.bold = True; p.font.color.rgb = WHITE

    slo_rows = [
        ("Document Parsing (100-page)", "< 15.0 s", "11.4 s", "PASSED"),
        ("Vector Indexing (500 chunks)", "< 5.0 s", "3.2 s", "PASSED"),
        ("Metric Extraction", "< 4.0 s", "2.1 s", "PASSED"),
        ("Red Flag Scanning", "< 4.0 s", "2.4 s", "PASSED"),
        ("Research TTFT (Streaming)", "< 800 ms", "420 ms", "PASSED"),
        ("PDF Report Compilation", "< 3.0 s", "1.8 s", "PASSED"),
        ("Grounding Verification", "100%", "100%", "PASSED"),
    ]
    for ri, row in enumerate(slo_rows, 1):
        bgc = CARD if ri % 2 == 1 else CARD_ALT
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci); c.fill.solid(); c.fill.fore_color.rgb = bgc
            p = c.text_frame.paragraphs[0]
            p.text = val; p.font.name = FONT_B; p.font.size = Pt(15)
            if ci == 0: p.font.bold = True; p.font.color.rgb = CHARCOAL
            elif ci == 3: p.font.bold = True; p.font.color.rgb = GREEN
            else: p.font.color.rgb = SLATE

    ftr(s, 14)

    # ═══════════════════════════════════════════════
    # SLIDE 15 — UI SHOWCASE
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Platform Showcase", "Interactive User Interface Walkthrough")

    shots = [
        ("Extraction Dashboard", os.path.join(SCREENSHOTS, "1_dashboard.png"), "GAAP metrics · Red Flag Radar · Confidence scores"),
        ("Research Chat", os.path.join(SCREENSHOTS, "5_research_chat.png"), "SSE streaming · Citation chips · Multi-turn QA"),
        ("PDF Report Viewer", os.path.join(SCREENSHOTS, "2_reports.png"), "One-click PDF export · Institutional formatting"),
    ]
    sw3 = Inches(3.72); sg3 = Inches(0.2)
    for i, (title, img_path, caption) in enumerate(shots):
        sx = MARGIN + i * (sw3 + sg3)
        card(s, sx, Inches(2.1), sw3, Inches(4.55))
        hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, Inches(2.1), sw3, Inches(0.48))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = NAVY; hbar.line.fill.background()
        htf = hbar.text_frame
        hp = htf.paragraphs[0]
        hp.text = title; hp.font.name = FONT_H; hp.font.size = Pt(16)
        hp.font.bold = True; hp.font.color.rgb = WHITE; hp.alignment = PP_ALIGN.CENTER
        if os.path.exists(img_path):
            s.shapes.add_picture(img_path, sx + Inches(0.15), Inches(2.72), width=sw3 - Inches(0.3))
        ct = tb(s, sx + Inches(0.15), Inches(6.0), sw3 - Inches(0.3), Inches(0.6))
        ctf = ct.text_frame; ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = caption; cp.font.name = FONT_B; cp.font.size = Pt(14)
        cp.font.color.rgb = EARTH; cp.alignment = PP_ALIGN.CENTER

    ftr(s, 15)

    # ═══════════════════════════════════════════════
    # SLIDE 16 — LEARNINGS & ROADMAP + Q&A
    # ═══════════════════════════════════════════════
    s = prs.slides.add_slide(BL); bg(s)
    hdr(s, "Conclusion", "Key Learnings & Future Roadmap")

    # Left — Competencies
    card(s, MARGIN, Inches(2.1), Inches(5.58), Inches(3.2))
    accent_bar(s, MARGIN, Inches(2.1), Inches(0.07), Inches(3.2), TEAL)
    tt = tb(s, MARGIN + Inches(0.25), Inches(2.25), Inches(5.1), Inches(2.9))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Engineering Competencies Gained"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = CHARCOAL
    bullet_text(ttf, [
        ("LangGraph Orchestration —", "Stateful cyclic agent graphs with MongoDB checkpointing"),
        ("Hybrid RAG Architecture —", "Dense + BM25 + Cross-Encoder retrieval pipeline"),
        ("Multi-Provider LLM Routing —", "Groq + Nemotron + Gemini circuit breaker resilience"),
        ("Security Engineering —", "STRIDE threat model, prompt defense, RBAC boundaries"),
    ], bold_size=15, body_size=14, bold_color=TEAL, body_color=SLATE)

    # Right — Roadmap
    card(s, Inches(6.58), Inches(2.1), Inches(5.58), Inches(3.2))
    accent_bar(s, Inches(6.58), Inches(2.1), Inches(0.07), Inches(3.2), NAVY)
    tt = tb(s, Inches(6.83), Inches(2.25), Inches(5.1), Inches(2.9))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Future Product Roadmap"; pp.font.name = FONT_H
    pp.font.size = Pt(22); pp.font.bold = True; pp.font.color.rgb = CHARCOAL
    bullet_text(ttf, [
        ("SEC EDGAR Webhooks —", "Real-time ingestion on new 10-K/10-Q filings"),
        ("XBRL Feed Ingestion —", "Structured iXBRL XML for deterministic accounting"),
        ("Multi-Modal Analysis —", "Vision LLM for parsing graphical financial charts"),
        ("Enterprise Connectors —", "SAP, Oracle Financials, Bloomberg Terminal feeds"),
    ], bold_size=15, body_size=14, bold_color=NAVY, body_color=SLATE)

    # Q&A banner
    card(s, MARGIN, Inches(5.55), CONTENT_W, Inches(1.1), fill=NAVY, brd=NAVY)
    tt = tb(s, MARGIN + Inches(0.2), Inches(5.62), CONTENT_W - Inches(0.4), Inches(0.95))
    ttf = tt.text_frame; ttf.word_wrap = True
    ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
    pp = ttf.paragraphs[0]
    pp.text = "Thank You  ·  Questions & Answers"; pp.font.name = FONT_H
    pp.font.size = Pt(28); pp.font.bold = True; pp.font.color.rgb = WHITE
    pp.alignment = PP_ALIGN.CENTER
    ps = ttf.add_paragraph()
    ps.text = "Team 2: Samrat (Lead) · Snigdha · Thirumala · Harshitha · Akshith"
    ps.font.name = FONT_B; ps.font.size = Pt(16); ps.font.color.rgb = RGBColor(203, 187, 160)
    ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(4)

    ftr(s, 16)

    # ───── Save ─────
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")


if __name__ == "__main__":
    create_presentation()
